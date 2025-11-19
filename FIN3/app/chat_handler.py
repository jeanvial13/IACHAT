import os
import time
import io
import json
from datetime import datetime, timedelta
from flask import (
    Flask,
    request,
    jsonify,
    render_template,
    send_file,
    redirect,
    url_for,
    url_for,
    session,
    Response,
    stream_with_context,
    send_from_directory,
)
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import base64
from werkzeug.utils import secure_filename
from openai import OpenAI

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    from docx import Document
except Exception:
    Document = None

try:
    from openpyxl import Workbook
except Exception:
    Workbook = None

try:
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.units import inch
except Exception:
    SimpleDocTemplate = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from docx.shared import Inches
except Exception:
    pass

BASE_DIR = os.path.dirname(os.path.abspath(__file__))

app = Flask(__name__, static_folder="static", template_folder="templates")

UPLOAD_FOLDER = os.path.join(BASE_DIR, "uploads")
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER

app.secret_key = os.environ.get("FLASK_SECRET_KEY", "dev-secret")

OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")
# Configuración
DEFAULT_MODEL = "gpt-4o"

client = OpenAI(api_key=OPENAI_API_KEY)

LOG_FILE = os.path.join(BASE_DIR, "server.log")
DEMS_FILE = os.path.join(BASE_DIR, "dem_projects.json")


# ---------------- Utilities ----------------


def _log(line: str) -> None:
    ts = datetime.utcnow().strftime("%Y-%m-%d %H:%M:%S")
    try:
        with open(LOG_FILE, "a", encoding="utf-8") as f:
            f.write(f"[{ts}] {line}\n")
    except Exception:
        pass


def extract_text(path: str) -> str:
    """Extract text from various file formats."""
    lower = path.lower()
    try:
        # Excel (.xlsx)
        if lower.endswith(".xlsx") and Workbook is not None:
            try:
                from openpyxl import load_workbook
                wb = load_workbook(path, data_only=True)
                text_parts = []
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    text_parts.append(f"--- Sheet: {sheet} ---")
                    for row in ws.iter_rows(values_only=True):
                        row_text = " | ".join([str(c) for c in row if c is not None])
                        if row_text.strip():
                            text_parts.append(row_text)
                return "\n".join(text_parts)
            except Exception as e:
                _log(f"Error reading Excel {path}: {e}")

        # PowerPoint (.pptx)
        if lower.endswith(".pptx") and Presentation is not None:
            try:
                prs = Presentation(path)
                text_parts = []
                for i, slide in enumerate(prs.slides):
                    text_parts.append(f"--- Slide {i+1} ---")
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_parts.append(shape.text)
                return "\n".join(text_parts)
            except Exception as e:
                _log(f"Error reading PPTX {path}: {e}")

        # PDF
        if lower.endswith(".pdf") and PdfReader is not None:
            reader = PdfReader(path)
            parts = []
            for page in reader.pages:
                try:
                    parts.append(page.extract_text() or "")
                except Exception:
                    pass
            return "\n".join(parts)

        # DOCX
        if lower.endswith(".docx") and Document is not None:
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs)

        # DOC (legacy)
        if lower.endswith(".doc"):
            try:
                import textract
                return textract.process(path).decode("utf-8", errors="ignore")
            except Exception:
                return ""
        
        # Try as plain text / code for everything else
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception:
            pass

    except Exception as e:
        _log(f"Error reading file {path}: {e}")

    return "Could not extract text from this file."


def load_dems():
    if not os.path.exists(DEMS_FILE):
        return []
    try:
        with open(DEMS_FILE, "r", encoding="utf-8") as f:
            data = f.read().strip()
            if not data:
                return []
            return json.loads(data)
    except Exception as e:
        _log(f"Error loading DEM file: {e}")
        return []


def save_dems(dems):
    try:
        with open(DEMS_FILE, "w", encoding="utf-8") as f:
            json.dump(dems, f, ensure_ascii=False, indent=2)
    except Exception as e:
        _log(f"Error saving DEM file: {e}")


def _clean_note_text(raw: str) -> str:
    """
    Quita una fecha duplicada al inicio si ya viene en el texto.

    Ejemplo:
        "[2025-11-18 00:52] — texto" -> "texto"
    """
    if not raw:
        return ""
    txt = raw.strip()
    if txt.startswith("[") and "] — " in txt:
        end = txt.find("] — ")
        if end != -1:
            return txt[end + 4 :].lstrip()
    return txt


def _format_note(note) -> str:
    """Devuelve una nota legible con la fecha, sin duplicarla."""
    if isinstance(note, dict):
        text = _clean_note_text(note.get("text", ""))
        date = note.get("date")
        if date:
            return f"[{date}] — {text}" if text else f"[{date}]"
        return text
    # compatibilidad con notas antiguas tipo string
    return _clean_note_text(str(note))


def enrich_dem(dem):
    """
    Agrega campos calculados:
    - duration_days
    - last_note
    - sla_breached
    - archived
    - priority
    - documents (lista)
    Además limpia el texto de las notas para quitar fechas duplicadas.
    """
    dem = dict(dem)

    # Normalizar notas (sin modificar el JSON en disco)
    raw_notes = dem.get("notes") or []
    cleaned_notes = []
    for n in raw_notes:
        if isinstance(n, dict):
            nn = dict(n)
            nn["text"] = _clean_note_text(nn.get("text", ""))
            cleaned_notes.append(nn)
        else:
            cleaned_notes.append(_clean_note_text(str(n)))
    dem["notes"] = cleaned_notes
    notes = cleaned_notes

    # Duración en días
    start_date = dem.get("start_date")
    dem["duration_days"] = None
    if start_date:
        try:
            dt = datetime.strptime(start_date, "%Y-%m-%d")
            dem["duration_days"] = (datetime.utcnow().date() - dt.date()).days
        except Exception:
            pass

    # Última nota formateada
    if notes:
        dem["last_note"] = _format_note(notes[-1])
    else:
        dem["last_note"] = ""

    # SLA (5 días sin actualización)
    updated_str = dem.get("updated_at") or dem.get("created_at")
    sla_breached = False
    if updated_str:
        try:
            upd = datetime.fromisoformat(updated_str)
            if datetime.utcnow() - upd > timedelta(days=5):
                sla_breached = True
        except Exception:
            pass
    dem["sla_breached"] = sla_breached

    # Archivado
    if "archived" not in dem:
        dem["archived"] = False

    # Prioridad por defecto
    if not dem.get("priority"):
        dem["priority"] = "2"

    # Lista de documentos
    docs = dem.get("documents")
    if docs is None or not isinstance(docs, list):
        dem["documents"] = []

    return dem


def generate_ai_comment(dem):
    """Generate a very short comment/highlight for the project using AI."""
    try:
        # Construct a prompt based on available data
        notes_text = "\n".join([n.get("text", "") if isinstance(n, dict) else str(n) for n in dem.get("notes", [])][-3:])
        prompt = (
            f"Project: {dem.get('name')}\n"
            f"Status: {dem.get('status')}\n"
            f"Recent Notes: {notes_text}\n\n"
            "Write a single, very short sentence (max 15 words) highlighting the most important thing about this project's current status or risk."
        )
        
        completion = client.chat.completions.create(
            model=DEFAULT_MODEL,
            messages=[
                {"role": "system", "content": "You are a helpful project manager assistant. Be concise."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=50,
        )
        return completion.choices[0].message.content.strip()
    except Exception as e:
        _log(f"Error generating AI comment: {e}")
        return "No AI comment available."


def build_portfolio_text(dems):
    """
    Construye el texto corporativo del portafolio para TXT/DOCX/PDF y el panel de UI.

    Sección 1: Projects Resume (resumen ejecutivo)
    Sección 2: Projects Details (detalle por DEM)

    IMPORTANTE:
      - Solo se deben pasar DEMs activos (no archivados).
    """
    if not dems:
        return "There are currently no DEM projects registered."

    enriched = [enrich_dem(d) for d in dems]

    run_date_human = datetime.utcnow().strftime("%B %d, %Y")
    header = f"{run_date_human} — Andres Villanueva DEMS Report"

    lines = []
    lines.append(header)
    lines.append("")
    lines.append("1. Projects Resume — Executive Overview")
    lines.append("")

    total = len(enriched)

    # Distribución de prioridades
    p1 = p2 = p3 = p4 = 0
    sla_ok = sla_breached = 0
    status_counts = {}

    for e in enriched:
        pr = str(e.get("priority") or "2")
        if pr == "1":
            p1 += 1
        elif pr == "2":
            p2 += 1
        elif pr == "3":
            p3 += 1
        elif pr == "4":
            p4 += 1

        if e.get("sla_breached"):
            sla_breached += 1
        else:
            sla_ok += 1

        st = e.get("status") or "N/A"
        status_counts[st] = status_counts.get(st, 0) + 1

    # Métricas clave
    lines.append("Key portfolio metrics for active DEM projects:")
    lines.append(f"• Total active DEMs: {total}")
    lines.append("• Priority distribution:")
    lines.append(f"   – P1 (Critical): {p1}")
    lines.append(f"   – P2 (High): {p2}")
    lines.append(f"   – P3 (Medium): {p3}")
    lines.append(f"   – P4 (Low): {p4}")
    lines.append(f"• SLA window (last 5 days): OK={sla_ok} | Breached={sla_breached}")

    if status_counts:
        top_status = sorted(status_counts.items(), key=lambda kv: kv[1], reverse=True)[
            :3
        ]
        status_str = ", ".join(f"{name}: {cnt}" for name, cnt in top_status)
        lines.append(f"• Most common DEM Status: {status_str}")

    lines.append("")
    lines.append("Active DEM overview (project name + latest comment):")
    lines.append("")

    for e in enriched:
        name = e.get("name", "(no name)")
        status = e.get("status", "-")
        workflow = e.get("workflow_status", "-")
        pr = e.get("priority", "2")

        raw_notes = e.get("notes") or []
        if raw_notes:
            latest = _format_note(raw_notes[-1])
        else:
            latest = "No recent notes registered."

        lines.append(
            f"• {name} — Status: {status} | Workflow: {workflow} | Priority: P{pr}"
        )
        lines.append(f"  Last update: {latest}")
        lines.append("")

    # Línea de separación donde marcaste en rojo
    lines.append("-" * 78)
    lines.append("")
    lines.append(
        "The following pages contain a detailed section per DEM, including "
        "Project Title, Sponsor, BA Owner, Workflow Status, SLA condition and "
        "the most recent notes captured during project follow-up."
    )
    lines.append("")
    lines.append("-" * 78)
    lines.append("")
    lines.append("2. Projects Details")
    lines.append("")

    # Detalle por DEM
    for e in enriched:
        lines.append(f"DEM: {e.get('name', '(no name)')}")
        lines.append(f"Project Title: {e.get('title', '')}")
        lines.append(
            f"Sponsor: {e.get('sponsor', '-')}"
            f" | Requester: {e.get('requester', '-')}"
        )
        lines.append(
            f"BA Owner: {e.get('ba_owner', '-')}"
            f" | Current Task Owner: {e.get('current_owner', '-')}"
        )
        lines.append(f"Cost Center: {e.get('cost_center', '-')}")
        lines.append(
            f"Start Date: {e.get('start_date', '-')}"
            f" | Duration (days): {e.get('duration_days')}"
        )
        lines.append(f"DEM Status: {e.get('status', '-')}")
        lines.append(f"Workflow Status: {e.get('workflow_status', '-')}")
        lines.append(f"Priority (1–4): {e.get('priority', '2')}")
        lines.append(
            "SLA Status: "
            + (
                "SLA Breached — project requires immediate follow-up with Sponsor and IT lead."
                if e.get("sla_breached")
                else "SLA OK — project updated within acceptable window."
            )
        )

        raw_notes = e.get("notes") or []
        if raw_notes:
            formatted_notes = [_format_note(n) for n in raw_notes]
            last_two = formatted_notes[-2:]
            lines.append("Last Notes (most recent entries):")
            for n in last_two:
                lines.append(f"- {n}")
        else:
            lines.append("Last Notes: (no notes registered)")

        lines.append("")
        lines.append("-" * 78)
        lines.append("")

    return "\n".join(lines)


# ---------------- Auth & Views ----------------


@app.route("/login", methods=["GET", "POST"])
def login():
    error = None
    if request.method == "POST":
        username = (request.form.get("username") or "").strip()
        password = request.form.get("password") or ""
        app_user = os.environ.get("APP_USER")
        app_pass = os.environ.get("APP_PASS")

        if app_user and app_pass and username == app_user and password == app_pass:
            session["auth"] = True
            return redirect(url_for("home"))
        error = "Usuario o contraseña incorrectos."
    return render_template("login.html", error=error)


@app.route("/logout")
def logout():
    session.clear()
    return redirect(url_for("login"))


def require_auth():
    if not session.get("auth"):
        return redirect(url_for("login"))
    return None


@app.route("/")
def home():
    maybe = require_auth()
    if maybe is not None:
        return maybe
    return render_template("index.html")


@app.route("/dems")
def dems_page():
    maybe = require_auth()
    if maybe is not None:
        return maybe
    return render_template("dem_manager.html")


# ---------------- Chat & Files (main chat page) ----------------


@app.route("/chat", methods=["POST"])
def chat():
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "No autorizado"}), 401

    data = request.get_json() or {}
    message = data.get("message", "").strip()
    history = data.get("history", [])
    model = data.get("model") or DEFAULT_MODEL
    file_summaries = data.get("file_summaries", [])

    if not message:
        return jsonify({"error": "Mensaje vacío"}), 400

    messages = []
    for m in history:
        role = m.get("role")
        content = m.get("content")
        if role in ("user", "assistant") and isinstance(content, str):
            messages.append({"role": role, "content": content})

    user_content = message

    if file_summaries:
        user_content += "\n\n[Attached files summary]\n"
        for fsum in file_summaries:
            fname = fsum.get("filename", "file")
            summ = fsum.get("summary", "")
            user_content += f"- {fname}: {summ}\n"

    messages.append({"role": "user", "content": user_content})

    def generate():
        try:
            stream = client.chat.completions.create(
                model=model,
                messages=messages,
                stream=True,
            )
            for chunk in stream:
                if chunk.choices[0].delta.content is not None:
                    yield chunk.choices[0].delta.content
        except Exception as e:
            _log(f"Error in OpenAI chat: {e}")
            yield f"Error: {str(e)}"

    return Response(stream_with_context(generate()), mimetype="text/plain")


@app.route("/upload", methods=["POST"])
def upload():
    """Upload generic files from the main chat and return short summaries."""
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "No autorizado"}), 401

    if "files" not in request.files:
        return jsonify({"error": "No files were sent."}), 400

    files = request.files.getlist("files")
    if not files:
        return jsonify({"error": "No files were sent."}), 400

    results = []
    for f in files:
        filename = secure_filename(f.filename or "file")
        save_name = f"{int(time.time())}_{filename}"
        path = os.path.join(app.config["UPLOAD_FOLDER"], save_name)
        _log(f"Saving uploaded file at {path}")

        f.save(path)

        text = extract_text(path)
        if not text:
            results.append(
                {
                    "filename": filename,
                    "summary": "I could not read this file (unsupported or empty).",
                }
            )
            continue

        try:
            completion = client.chat.completions.create(
                model=DEFAULT_MODEL,
                messages=[
                    {
                        "role": "user",
                        "content": (
                            "Summarize the following document in a few bullet points, "
                            "highlighting key information useful for IT, business analysis "
                            "and project follow-up:\n\n"
                            f"{text[:8000]}"
                        ),
                    }
                ],
            )
            summary = completion.choices[0].message.content
        except Exception as e:
            _log(f"Error summarizing file: {e}")
            summary = (
                "An automatic summary could not be generated, "
                "but the file was uploaded correctly."
            )

        results.append({"filename": filename, "summary": summary})

    return jsonify({"files": results})


# ---------------- DEM / Project Manager API ----------------


def get_dems_filtered(archived: bool):
    dems = load_dems()
    return [enrich_dem(d) for d in dems if bool(d.get("archived", False)) == archived]


@app.route("/api/dems/projects", methods=["GET"])
def list_dems():
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "No autorizado"}), 401

    archived_str = request.args.get("archived", "false").lower()
    archived = archived_str in ("1", "true", "yes")
    projects = get_dems_filtered(archived)
    return jsonify({"projects": projects})


@app.route("/api/dems/projects", methods=["POST"])
def create_dem():
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "No autorizado"}), 401

    data = request.get_json() or {}
    now_iso = datetime.utcnow().isoformat()

    dem = {
        "id": f"dem_{int(time.time() * 1000)}",
        "name": data.get("name", "").strip(),
        "sponsor": data.get("sponsor", "").strip(),
        "requester": data.get("requester", "").strip(),
        "ba_owner": data.get("ba_owner", "").strip(),
        "title": data.get("title", "").strip(),
        "change_request": data.get("change_request", "").strip(),
        "cost_center": data.get("cost_center", "").strip(),
        "status": data.get("status", "").strip() or "Idea",
        "workflow_status": data.get("workflow_status", "").strip() or "Intake",
        "current_owner": data.get("current_owner", "").strip(),
        "start_date": data.get("start_date", "").strip(),
        "priority": (data.get("priority") or "2").strip(),
        "notes": [],
        "documents": [],  # multiple documents per DEM
        "doc_summary": "",
        "created_at": now_iso,
        "updated_at": now_iso,
        "archived": False,
    }

    initial_note = data.get("initial_note", "").strip() if "initial_note" in data else ""
    if initial_note:
        dem["notes"].append(
            {
                "text": _clean_note_text(initial_note),
                "date": datetime.utcnow().strftime("%Y-%m-%d %H:%M"),
            }
        )

    dems = load_dems()
    dems.append(dem)
    save_dems(dems)

    return jsonify({"project": enrich_dem(dem)})


def _update_dem(id, updater):
    dems = load_dems()
    for i, d in enumerate(dems):
        if d.get("id") == id:
            updater(d)
            d["updated_at"] = datetime.utcnow().isoformat()
            dems[i] = d
            save_dems(dems)
            return enrich_dem(d)
    return None


@app.route("/api/dems/projects/<id>/note", methods=["POST"])
def add_dem_note(id):
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "No autorizado"}), 401

    data = request.get_json() or {}
    text = (data.get("text") or "").strip()
    if not text:
        return jsonify({"error": "Nota vacía."}), 400

    clean_text = _clean_note_text(text)

    def updater(d):
        notes = d.get("notes") or []
        if not isinstance(notes, list):
            notes = []
        notes.append(
            {
                "text": clean_text,
                "date": datetime.utcnow().strftime("%Y-%m-%d %H:%M"),
            }
        )
        d["notes"] = notes

    project = _update_dem(id, updater)
    if not project:
        return jsonify({"error": "DEM no encontrado."}), 404
    return jsonify({"project": project})


# ---- EDIT NOTE ---------------------------------------------------------
@app.route("/api/dems/projects/<id>/note/edit", methods=["POST"])
def edit_dem_note(id):
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "Unauthorized"}), 401

    data = request.get_json() or {}
    index = data.get("index")
    new_text = (data.get("text") or "").strip()

    if index is None or new_text == "":
        return jsonify({"error": "Invalid index or empty text."}), 400

    clean_text = _clean_note_text(new_text)

    def updater(d):
        notes = d.get("notes") or []
        if not isinstance(notes, list):
            notes = []
        if index < 0 or index >= len(notes):
            raise ValueError("Invalid index")

        notes[index] = {
            "text": clean_text,
            "date": datetime.utcnow().strftime("%Y-%m-%d %H:%M"),
        }
        d["notes"] = notes

    try:
        project = _update_dem(id, updater)
    except Exception as e:
        return jsonify({"error": str(e)}), 400

    if not project:
        return jsonify({"error": "DEM not found"}), 404
    return jsonify({"project": project})


# ---- DELETE NOTE ---------------------------------------------------------
@app.route("/api/dems/projects/<id>/note/delete", methods=["POST"])
def delete_dem_note(id):
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "Unauthorized"}), 401

    data = request.get_json() or {}
    index = data.get("index")

    if index is None:
        return jsonify({"error": "Missing index"}), 400

    def updater(d):
        notes = d.get("notes") or []
        if not isinstance(notes, list):
            notes = []

        if index < 0 or index >= len(notes):
            raise ValueError("Invalid index")

        notes.pop(index)
        d["notes"] = notes

    try:
        project = _update_dem(id, updater)
    except Exception as e:
        return jsonify({"error": str(e)}), 400

    if not project:
        return jsonify({"error": "DEM not found"}), 404

    return jsonify({"project": project})


@app.route("/api/dems/projects/<id>/update", methods=["POST"])
def update_dem(id):
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "No autorizado"}), 401

    data = request.get_json() or {}

    def updater(d):
        for field in [
            "name",
            "sponsor",
            "requester",
            "ba_owner",
            "title",
            "change_request",
            "cost_center",
            "status",
            "workflow_status",
            "current_owner",
            "start_date",
            "priority",
        ]:
            if field in data:
                d[field] = (data.get(field) or "").strip()

    project = _update_dem(id, updater)
    if not project:
        return jsonify({"error": "DEM no encontrado."}), 404
    return jsonify({"project": project})


@app.route("/api/dems/projects/<id>/archive", methods=["POST"])
def archive_dem(id):
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "No autorizado"}), 401

    def updater(d):
        d["archived"] = True

    project = _update_dem(id, updater)
    if not project:
        return jsonify({"error": "DEM no encontrado."}), 404
    return jsonify({"project": project})


@app.route("/api/dems/projects/<id>/restore", methods=["POST"])
def restore_dem(id):
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "No autorizado"}), 401

    def updater(d):
        d["archived"] = False

    project = _update_dem(id, updater)
    if not project:
        return jsonify({"error": "DEM no encontrado."}), 404
    return jsonify({"project": project})


@app.route("/api/dems/projects/<id>/delete", methods=["POST"])
def delete_dem(id):
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "No autorizado"}), 401

    dems = load_dems()
    new_dems = [d for d in dems if d.get("id") != id]
    if len(new_dems) == len(dems):
        return jsonify({"error": "DEM no encontrado."}), 404
    save_dems(new_dems)
    return jsonify({"success": True})


@app.route("/api/dems/projects/<id>/attach", methods=["POST"])
def attach_doc(id):
    """Attach a document to a DEM, analyze it and store the summary."""
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "No autorizado"}), 401

    if "file" not in request.files:
        return jsonify({"error": "No se recibió archivo."}), 400

    file = request.files["file"]
    filename = secure_filename(file.filename or "document")
    save_name = f"{int(time.time())}_{filename}"
    path = os.path.join(app.config["UPLOAD_FOLDER"], save_name)
    file.save(path)

    text = extract_text(path)
    if not text:
        return jsonify(
            {
                "error": "I could not read this file to generate an executive summary."
            }
        ), 400

    summary = ""
    try:
        completion = client.chat.completions.create(
            model=DEFAULT_MODEL,
            messages=[
                {
                    "role": "user",
                    "content": (
                        "Create an executive summary of the following document for a DEM "
                        "(Digital Enhancement Management) portfolio. "
                        "Focus on: business problem, scope, key requirements, risks, "
                        "dependencies, recommended IT solutions (for example SAP S/4HANA "
                        "or other enterprise systems), and clear next actions. "
                        "Write in concise, professional English. "
                        "Do NOT mention that this text was generated by any AI model and "
                        "do not describe any internal technical process.\n\n"
                        f"{text[:8000]}"
                    ),
                }
            ],
        )
        summary = completion.choices[0].message.content
    except Exception as e:
        _log(f"Error generating doc summary: {e}")
        summary = (
            "An automatic executive summary could not be generated, "
            "but the document was attached to this DEM."
        )

    def updater(d):
        # Ensure documents list exists
        docs = d.get("documents") or []
        if not isinstance(docs, list):
            docs = []
        docs.append(
            {
                "filename": filename,
                "summary": summary,
                "date": datetime.utcnow().strftime("%Y-%m-%d %H:%M"),
            }
        )
        d["documents"] = docs
        # Keep last summary in doc_summary for current UI button
        d["doc_summary"] = summary

    project = _update_dem(id, updater)
    if not project:
        return jsonify({"error": "DEM no encontrado."}), 404
    return jsonify({"project": project})


@app.route("/api/dems/projects/<id>/summary/delete", methods=["POST"])
def delete_dem_summary(id):
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "Unauthorized"}), 401

    def updater(d):
        d["doc_summary"] = ""

    project = _update_dem(id, updater)
    if not project:
        return jsonify({"error": "DEM not found"}), 404
    return jsonify({"project": project})


@app.route("/api/dems/export", methods=["GET"])
def export_active_excel():
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "No autorizado"}), 401

    if Workbook is None:
        return jsonify({"error": "openpyxl no está disponible en el servidor."}), 500

    wb = Workbook()
    ws = wb.active
    ws.title = "Active DEMs"

    headers = [
        "ID",
        "Name",
        "Title",
        "Sponsor",
        "Requester",
        "BA Owner",
        "Cost Center",
        "Status",
        "Workflow Status",
        "Current Task Owner",
        "Start Date",
        "Duration Days",
        "SLA",
        "Last Note",
    ]
    ws.append(headers)

    for dem in get_dems_filtered(False):
        ws.append(
            [
                dem.get("id"),
                dem.get("name"),
                dem.get("title"),
                dem.get("sponsor"),
                dem.get("requester"),
                dem.get("ba_owner"),
                dem.get("cost_center"),
                dem.get("status"),
                dem.get("workflow_status"),
                dem.get("current_owner"),
                dem.get("start_date"),
                dem.get("duration_days"),
                "Breached" if dem.get("sla_breached") else "OK",
                dem.get("last_note"),
            ]
        )

    bio = io.BytesIO()
    wb.save(bio)
    bio.seek(0)
    return send_file(
        bio,
        as_attachment=True,
        download_name="dems_active.xlsx",
        mimetype="application/"
        "vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )


@app.route("/api/dems/export_archived", methods=["GET"])
def export_archived_excel():
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "No autorizado"}), 401

    if Workbook is None:
        return jsonify({"error": "openpyxl no está disponible en el servidor."}), 500

    wb = Workbook()
    ws = wb.active
    ws.title = "Archived DEMs"

    headers = [
        "ID",
        "Name",
        "Title",
        "Sponsor",
        "Requester",
        "BA Owner",
        "Cost Center",
        "Status",
        "Workflow Status",
        "Current Task Owner",
        "Start Date",
        "Duration Days",
        "SLA",
        "Last Note",
    ]
    ws.append(headers)

    for dem in get_dems_filtered(True):
        ws.append(
            [
                dem.get("id"),
                dem.get("name"),
                dem.get("title"),
                dem.get("sponsor"),
                dem.get("requester"),
                dem.get("ba_owner"),
                dem.get("cost_center"),
                dem.get("status"),
                dem.get("workflow_status"),
                dem.get("current_owner"),
                dem.get("start_date"),
                dem.get("duration_days"),
                "Breached" if dem.get("sla_breached") else "OK",
                dem.get("last_note"),
            ]
        )

    bio = io.BytesIO()
    wb.save(bio)
    bio.seek(0)
    return send_file(
        bio,
        as_attachment=True,
        download_name="dems_archived.xlsx",
        mimetype="application/"
        "vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )


# -------- Export / Import JSON (backup) ----------------------


@app.route("/api/dems/export_json", methods=["GET"])
def export_dems_json():
    """
    Exporta TODOS los DEMs (activos + archivados) como backup JSON.
    El botón “Export JSON” del frontend debe llamar a este endpoint.
    """
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "No autorizado"}), 401

    dems = load_dems() or []
    payload = json.dumps(dems, ensure_ascii=False, indent=2)
    bio = io.BytesIO(payload.encode("utf-8"))
    bio.seek(0)
    return send_file(
        bio,
        as_attachment=True,
        download_name="dems_backup.json",
        mimetype="application/json; charset=utf-8",
    )


@app.route("/api/dems/import", methods=["POST"])
def import_dems_json():
    """
    Importa DEMs desde un JSON.

    El frontend normalmente manda algo así:
        POST /api/dems/import
        { "projects": [ {..dem1..}, {..dem2..}, ... ] }
    """
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "No autorizado"}), 401

    data = request.get_json(silent=True) or {}
    projects = data.get("projects")

    if not isinstance(projects, list):
        return jsonify(
            {"error": "Invalid JSON structure: 'projects' must be a list."}
        ), 400

    current = load_dems() or []
    by_id = {str(d.get("id")): d for d in current if d.get("id")}

    for incoming in projects:
        if not isinstance(incoming, dict):
            continue
        pid = str(incoming.get("id") or "").strip()
        if not pid:
            pid = f"dem_{int(time.time() * 1000)}"
            incoming["id"] = pid
        by_id[pid] = incoming

    merged_list = list(by_id.values())
    save_dems(merged_list)

    enriched = [enrich_dem(d) for d in merged_list]
    return jsonify({"projects": enriched})


# -------- Reporte resumen para panel y descargas -------------


@app.route("/api/dems/report", methods=["POST"])
def dem_report():
    """Texto del reporte para el panel (solo DEMs activos)."""
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "No autorizado"}), 401

    dems = load_dems()
    active_dems = [d for d in dems if not d.get("archived", False)]
    # text = build_portfolio_text(active_dems)
    html_report = build_portfolio_html(active_dems)

    return jsonify({"report": html_report})


def build_portfolio_html(projects):
    """Generates a Cyberpunk-styled HTML report for the portfolio."""
    now_str = datetime.utcnow().strftime("%Y-%m-%d %H:%M")
    
    html = f"""
    <div style="font-family: 'Rajdhani', sans-serif; color: #e0e0e0;">
        <div style="border-bottom: 2px solid #00f3ff; padding-bottom: 20px; margin-bottom: 30px;">
            <h1 style="color: #fcee0a; font-size: 32px; margin: 0; text-transform: uppercase; letter-spacing: 2px;">
                Andres Villanueva DEMS Report
            </h1>
            <div style="display: flex; justify-content: space-between; margin-top: 10px; font-family: 'Roboto Mono', monospace; font-size: 12px; color: #00f3ff;">
                <span>GENERATED: {now_str}</span>
                <span>CONFIDENTIAL // INTERNAL USE ONLY</span>
            </div>
        </div>

        <div style="margin-bottom: 40px;">
            <h2 style="color: #ff00ff; border-left: 4px solid #ff00ff; padding-left: 15px; margin-bottom: 20px; text-transform: uppercase;">
                1. Portfolio Dashboard
            </h2>
            <div style="display: grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr)); gap: 20px;">
                <div style="background: rgba(0, 243, 255, 0.05); border: 1px solid #00f3ff; padding: 20px;">
                    <div style="font-size: 12px; color: #888899; text-transform: uppercase;">Active Projects</div>
                    <div style="font-size: 36px; color: #fff; font-weight: bold;">{len(projects)}</div>
                </div>
                <div style="background: rgba(255, 42, 42, 0.05); border: 1px solid #ff2a2a; padding: 20px;">
                    <div style="font-size: 12px; color: #888899; text-transform: uppercase;">Priority 1 (Critical)</div>
                    <div style="font-size: 36px; color: #ff2a2a; font-weight: bold;">
                        {len([p for p in projects if str(p.get('priority')) == '1'])}
                    </div>
                </div>
                 <div style="background: rgba(252, 238, 10, 0.05); border: 1px solid #fcee0a; padding: 20px;">
                    <div style="font-size: 12px; color: #888899; text-transform: uppercase;">SLA Breached</div>
                    <div style="font-size: 36px; color: #fcee0a; font-weight: bold;">
                        {len([p for p in projects if p.get('sla_breached')])}
                    </div>
                </div>
            </div>
        </div>

        <div>
            <h2 style="color: #00f3ff; border-left: 4px solid #00f3ff; padding-left: 15px; margin-bottom: 20px; text-transform: uppercase;">
                2. Project Details
            </h2>
    """

    for p in projects:
        p_id = p.get("id", "N/A")
        name = p.get("name", "Untitled")
        status = p.get("status", "Unknown")
        workflow = p.get("workflow_status", "Unknown")
        priority = p.get("priority", "4")
        owner = p.get("current_owner", "Unassigned")
        notes = p.get("notes", [])
        last_note = notes[-1]["text"] if notes else "No notes available."
        last_date = notes[-1]["date"] if notes else "N/A"
        
        priority_color = "#888"
        if str(priority) == "1": priority_color = "#ff2a2a"
        elif str(priority) == "2": priority_color = "#ffaa00"
        elif str(priority) == "3": priority_color = "#00f3ff"

        html += f"""
        <div style="background: rgba(10, 10, 15, 0.8); border: 1px solid #333; margin-bottom: 20px; padding: 20px; position: relative;">
            <div style="position: absolute; top: 0; left: 0; width: 4px; height: 100%; background: {priority_color};"></div>
            <div style="display: flex; justify-content: space-between; align-items: flex-start; margin-bottom: 15px; padding-left: 15px;">
                <div>
                    <h3 style="margin: 0; font-size: 18px; color: #fff;">{name}</h3>
                    <div style="font-family: 'Roboto Mono', monospace; font-size: 11px; color: #888899; margin-top: 4px;">
                        ID: {p_id} | OWNER: {owner}
                    </div>
                </div>
                <div style="text-align: right;">
                    <span style="background: {priority_color}20; color: {priority_color}; padding: 4px 8px; border: 1px solid {priority_color}; font-size: 10px; font-family: 'Roboto Mono', monospace;">
                        P{priority}
                    </span>
                </div>
            </div>
            
            <div style="display: grid; grid-template-columns: 1fr 1fr; gap: 15px; margin-bottom: 15px; padding-left: 15px; font-family: 'Roboto Mono', monospace; font-size: 12px;">
                <div><span style="color: #00f3ff;">STATUS:</span> {status}</div>
                <div><span style="color: #00f3ff;">WORKFLOW:</span> {workflow}</div>
            </div>

            <div style="background: rgba(0, 0, 0, 0.3); padding: 15px; border-left: 2px solid #555; margin-left: 15px;">
                <div style="font-family: 'Roboto Mono', monospace; font-size: 10px; color: #00f3ff; margin-bottom: 5px;">
                    LATEST UPDATE [{last_date}]
                </div>
                <div style="font-size: 13px; line-height: 1.5; color: #ccc; margin-bottom: 15px;">
                    {last_note}
                </div>

                <!-- AI Comment Section -->
                <div style="border-top: 1px dashed #333; padding-top: 10px; margin-top: 10px;">
                    <div style="font-family: 'Roboto Mono', monospace; font-size: 10px; color: #fcee0a; margin-bottom: 5px; text-transform: uppercase;">
                        AI Insight // Project Summary
                    </div>
                    <div style="font-size: 13px; line-height: 1.4; color: #e0e0e0; font-style: italic;">
                        "{generate_ai_comment(p)}"
                    </div>
                </div>
            </div>
        </div>
        """

    html += """
        <div style="margin-top: 50px; border-top: 1px solid #333; padding-top: 20px; text-align: center; font-family: 'Roboto Mono', monospace; font-size: 10px; color: #555;">
            END OF REPORT // ANDRES IA SYSTEM
        </div>
    </div>
    """
    
    return html


def generate_charts(dems):
    """Generate pie and bar charts for the report."""
    charts = {}
    
    # 1. Priority Distribution (Horizontal Bar)
    priorities = {"1": 0, "2": 0, "3": 0, "4": 0}
    for d in dems:
        p = str(d.get("priority") or "2")
        if p in priorities:
            priorities[p] += 1
    
    # Sort by priority key (1, 2, 3, 4)
    sorted_keys = sorted(priorities.keys(), reverse=True) # P4 top, P1 bottom? Or P1 top? Barh usually goes bottom-up.
    # Let's do P1 at top. So reverse=True for keys if we plot 1,2,3,4.
    # Actually, let's just list them P1, P2, P3, P4.
    keys = ["1", "2", "3", "4"]
    values = [priorities[k] for k in keys]
    labels = [f"P{k}" for k in keys]
    
    # Neon colors: Red, Amber, Blue, Emerald
    colors_prio = ['#ff2a2a', '#fcee0a', '#00f3ff', '#00ff9f']
    
    fig1, ax1 = plt.subplots(figsize=(5, 3))
    # Horizontal bar
    y_pos = range(len(keys))
    ax1.barh(y_pos, values, color=colors_prio, height=0.6)
    ax1.set_yticks(y_pos)
    ax1.set_yticklabels(labels, color="white")
    ax1.invert_yaxis() # P1 at top
    
    # Add value labels
    for i, v in enumerate(values):
        ax1.text(v + 0.1, i, str(v), color='white', va='center', fontweight='bold')

    ax1.set_title("Priority Distribution", color="white", pad=10)
    # Remove borders/ticks for cleaner look
    ax1.spines['top'].set_visible(False)
    ax1.spines['right'].set_visible(False)
    ax1.spines['bottom'].set_visible(False)
    ax1.spines['left'].set_visible(False)
    ax1.tick_params(axis='x', colors='white')
    ax1.tick_params(axis='y', colors='white')
    
    # Transparent bg
    fig1.patch.set_alpha(0.0)
    ax1.patch.set_alpha(0.0)

    buf1 = io.BytesIO()
    plt.savefig(buf1, format='png', transparent=True, dpi=150, bbox_inches='tight')
    buf1.seek(0)
    charts['priority'] = buf1
    plt.close(fig1)

    # 2. SLA Status (Horizontal Bar)
    sla_counts = {"OK": 0, "Breached": 0}
    for d in dems:
        if d.get("sla_breached"):
            sla_counts["Breached"] += 1
        else:
            sla_counts["OK"] += 1
            
    fig2, ax2 = plt.subplots(figsize=(5, 2))
    # Emerald for OK, Red for Breached
    keys_sla = ["OK", "Breached"]
    vals_sla = [sla_counts["OK"], sla_counts["Breached"]]
    colors_sla = ['#00ff9f', '#ff2a2a']
    
    y_pos2 = range(len(keys_sla))
    ax2.barh(y_pos2, vals_sla, color=colors_sla, height=0.5)
    ax2.set_yticks(y_pos2)
    ax2.set_yticklabels(keys_sla, color="white")
    ax2.invert_yaxis()
    
    for i, v in enumerate(vals_sla):
        ax2.text(v + 0.1, i, str(v), color='white', va='center', fontweight='bold')

    ax2.set_title("SLA Status", color="white", pad=10)
    ax2.spines['top'].set_visible(False)
    ax2.spines['right'].set_visible(False)
    ax2.spines['bottom'].set_visible(False)
    ax2.spines['left'].set_visible(False)
    ax2.tick_params(axis='x', colors='white')
    ax2.tick_params(axis='y', colors='white')
    
    fig2.patch.set_alpha(0.0)
    ax2.patch.set_alpha(0.0)
    
    buf2 = io.BytesIO()
    plt.savefig(buf2, format='png', transparent=True, dpi=150, bbox_inches='tight')
    buf2.seek(0)
    charts['sla'] = buf2
    plt.close(fig2)

    # 3. Project Status (Horizontal Bar)
    status_counts = {}
    for d in dems:
        s = d.get("status", "Unknown")
        status_counts[s] = status_counts.get(s, 0) + 1
    
    # Sort by count
    sorted_status = sorted(status_counts.items(), key=lambda x: x[1], reverse=True)
    keys_stat = [x[0] for x in sorted_status]
    vals_stat = [x[1] for x in sorted_status]
    
    fig3, ax3 = plt.subplots(figsize=(6, len(keys_stat)*0.5 + 1))
    # Neon palette
    colors_status = ['#00f3ff', '#bd00ff', '#fcee0a', '#ff2a2a', '#00ff9f', '#ffffff']
    # Cycle colors
    bar_colors = [colors_status[i % len(colors_status)] for i in range(len(keys_stat))]
    
    y_pos3 = range(len(keys_stat))
    ax3.barh(y_pos3, vals_stat, color=bar_colors, height=0.6)
    ax3.set_yticks(y_pos3)
    ax3.set_yticklabels(keys_stat, color="white", fontsize=8)
    ax3.invert_yaxis()
    
    for i, v in enumerate(vals_stat):
        ax3.text(v + 0.1, i, str(v), color='white', va='center', fontweight='bold')

    ax3.set_title("Project Status", color="white", pad=10)
    ax3.spines['top'].set_visible(False)
    ax3.spines['right'].set_visible(False)
    ax3.spines['bottom'].set_visible(False)
    ax3.spines['left'].set_visible(False)
    ax3.tick_params(axis='x', colors='white')
    ax3.tick_params(axis='y', colors='white')

    fig3.patch.set_alpha(0.0)
    ax3.patch.set_alpha(0.0)

    buf3 = io.BytesIO()
    plt.savefig(buf3, format='png', transparent=True, dpi=150, bbox_inches='tight')
    buf3.seek(0)
    charts['status'] = buf3
    plt.close(fig3)
    
    return charts

@app.route("/api/dems/download/<fmt>", methods=["GET"])
def dem_download(fmt):
    """Descarga el reporte como TXT / DOCX / PDF (solo DEMs activos)."""
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "No autorizado"}), 401

    fmt = fmt.lower()
    if fmt not in ("txt", "pdf", "docx"):
        return jsonify({"error": "Formato no soportado."}), 400

    dems = [d for d in load_dems() if not d.get("archived", False)]
    text = build_portfolio_text(dems)
    charts = generate_charts(dems)

    lines = text.split("\n")
    title = lines[0] if lines else "Andres Villanueva DEMS Report"
    
    # TXT (Simple dump)
    if fmt == "txt":
        bio = io.BytesIO()
        bio.write(text.encode("utf-8"))
        bio.seek(0)
        return send_file(
            bio,
            as_attachment=True,
            download_name="dems_portfolio.txt",
            mimetype="text/plain; charset=utf-8",
        )

    # DOCX (Enhanced)
    if fmt == "docx":
        if Document is None:
            return jsonify({"error": "python-docx no está disponible."}), 500
        doc = Document()
        doc.add_heading(title, level=0)
        doc.add_paragraph(f"Generated on: {datetime.utcnow().strftime('%Y-%m-%d %H:%M')}")
        
        # Dashboard Section
        doc.add_heading("1. Portfolio Dashboard", level=1)
        
        # Add charts side by side (using a table)
        table = doc.add_table(rows=1, cols=2)
        table.autofit = True
        
        # Priority Chart
        cell1 = table.cell(0, 0)
        p1 = cell1.paragraphs[0]
        run1 = p1.add_run()
        run1.add_picture(charts['priority'], width=Inches(2.8))
        
        # SLA Chart
        cell2 = table.cell(0, 1)
        p2 = cell2.paragraphs[0]
        run2 = p2.add_run()
        run2.add_picture(charts['sla'], width=Inches(2.8))
        
        doc.add_paragraph("") # Spacer

        # Text Content
        doc.add_heading("2. Project Details", level=1)
        
        # Parse the text report and add it nicely
        # Skipping the first few lines which are the header/summary in the text version
        start_details = False
        for line in lines:
            if "2. Projects Details" in line:
                start_details = True
                continue
            if not start_details:
                if line.strip() and "1. Projects Resume" not in line and title not in line:
                     doc.add_paragraph(line)
                continue
            
            if line.startswith("DEM:"):
                doc.add_heading(line, level=2)
            elif line.startswith("Last Notes"):
                doc.add_heading(line, level=3)
            elif line.strip() == "-" * 78:
                pass # Skip separators
            else:
                doc.add_paragraph(line)

        bio = io.BytesIO()
        doc.save(bio)
        bio.seek(0)
        return send_file(
            bio,
            as_attachment=True,
            download_name="dems_portfolio.docx",
            mimetype=(
                "application/"
                "vnd.openxmlformats-officedocument.wordprocessingml.document"
            ),
        )

    # PDF (Enhanced Dashboard)
    if fmt == "pdf":
        if SimpleDocTemplate is None:
            return jsonify({"error": "reportlab no está disponible."}), 500

        bio = io.BytesIO()
        doc = SimpleDocTemplate(bio, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        
        # Title: Andres Villanueva Project Status
        title_style = styles["Title"]
        title_style.textColor = colors.HexColor("#3b82f6")
        story.append(Paragraph("Andres Villanueva Project Status", title_style))
        story.append(Spacer(1, 6))
        story.append(Paragraph(f"Generated on: {datetime.utcnow().strftime('%Y-%m-%d %H:%M')}", styles["Normal"]))
        story.append(Spacer(1, 24))
        
        # 1. Active Projects Count
        story.append(Paragraph("Active Projects", styles["Heading2"]))
        active_count = len(dems)
        # Create a big number style
        big_num_style = ParagraphStyle('BigNum', parent=styles['Normal'], fontSize=36, leading=42, textColor=colors.HexColor("#10b981"), alignment=1) # Center
        story.append(Paragraph(str(active_count), big_num_style))
        story.append(Spacer(1, 24))

        # 2. Charts (Priority, Status)
        story.append(Paragraph("Portfolio Metrics", styles["Heading2"]))
        story.append(Spacer(1, 12))
        
        # Row of charts
        chart_row = [
            Image(charts['priority'], width=2.5*inch, height=2*inch),
            Image(charts['status'], width=3*inch, height=2*inch)
        ]
        t_charts = Table([chart_row])
        t_charts.setStyle(TableStyle([
            ('ALIGN', (0,0), (-1,-1), 'CENTER'),
            ('VALIGN', (0,0), (-1,-1), 'MIDDLE')
        ]))
        story.append(t_charts)
        story.append(Spacer(1, 24))

        # 3. Last Comments Dashboard
        story.append(Paragraph("Last Comments / Status Updates", styles["Heading2"]))
        story.append(Spacer(1, 12))

        # Build table data for comments
        # Header
        table_data = [[Paragraph("<b>Project</b>", styles["Normal"]), Paragraph("<b>Last Comment</b>", styles["Normal"])]]
        
        for d in dems:
            name = d.get("name", "Unknown")
            notes = d.get("notes", [])
            
            if notes:
                last_note_obj = notes[-1]
                last_note_text = last_note_obj.get("text", "")
                last_note_date = last_note_obj.get("date", "N/A")
                # Combine date and text
                full_comment = f"<b>[{last_note_date}]</b> {last_note_text}"
            else:
                full_comment = "No notes available."

            # Truncate if too long for table
            if len(full_comment) > 300:
                full_comment = full_comment[:300] + "..."
            
            table_data.append([
                Paragraph(name, styles["Normal"]),
                Paragraph(full_comment, ParagraphStyle('MonoComment', parent=styles["Normal"], fontName='Courier', fontSize=9, backColor=colors.HexColor('#f0f0f0'), borderPadding=4))
            ])

        t_comments = Table(table_data, colWidths=[2*inch, 4.5*inch])
        t_comments.setStyle(TableStyle([
            ('GRID', (0,0), (-1,-1), 0.5, colors.grey),
            ('BACKGROUND', (0,0), (1,0), colors.lightgrey),
            ('VALIGN', (0,0), (-1,-1), 'TOP'),
            ('PADDING', (0,0), (-1,-1), 6),
        ]))
        story.append(t_comments)
        story.append(Spacer(1, 24))
        
        # 4. Detailed Project List (Standard Report)
        story.append(Paragraph("Detailed Project Report", styles["Heading1"]))
        story.append(Spacer(1, 12))
        
        # Parse text for PDF details
        start_details = False
        for line in lines:
            if "2. Projects Details" in line:
                start_details = True
                continue
            
            if not start_details:
                # Skip the old summary part since we have a new dashboard
                continue

            if line.startswith("DEM:"):
                story.append(Spacer(1, 12))
                story.append(Paragraph(line, styles["Heading2"]))
            elif line.startswith("Last Notes"):
                story.append(Paragraph(line, styles["Heading3"]))
            elif line.strip() == "-" * 78:
                story.append(Spacer(1, 6))
                story.append(Paragraph("_" * 60, styles["BodyText"]))
                story.append(Spacer(1, 6))
            else:
                # Handle bullet points
                if line.strip().startswith("- ") or line.strip().startswith("• "):
                    story.append(Paragraph(line, styles["Bullet"], bulletText="•"))
                else:
                    story.append(Paragraph(line, styles["Normal"]))

        doc.build(story)
        bio.seek(0)
        return send_file(
            bio,
            as_attachment=True,
            download_name="dems_portfolio.pdf",
            mimetype="application/pdf",
        )

    return jsonify({"error": "Formato no soportado."}), 400


# ---------------- AMD AI REPORTS ----------------

def generate_amd_ai_report_logic(projects):
    """
    Generates a strategic report using OpenAI based on the provided projects.
    Now includes document summaries in the context.
    """
    if not projects:
        return "No active projects to analyze."

    # Build a context string from the projects
    project_context = ""
    for p in projects:
        name = p.get("name", "Unknown")
        title = p.get("title", "No Title")
        status = p.get("status", "Unknown")
        workflow = p.get("workflow_status", "Unknown")
        priority = p.get("priority", "N/A")
        
        # Get document summary if available
        doc_summary = p.get("doc_summary", "No document summary available.")
        
        project_context += f"- Project: {name} | Title: {title}\n"
        project_context += f"  Status: {status} | Workflow: {workflow} | Priority: {priority}\n"
        project_context += f"  Document Summary: {doc_summary}\n"
        project_context += "  ---\n"

    system_prompt = (
        "You are a Senior IT Strategic Advisor for AMD, specializing in SAP S/4HANA, Cloud Migrations, and Enterprise Architecture. "
        "Your goal is to provide a high-level executive summary and actionable strategic advice for the following portfolio of projects.\n\n"
        "Focus on providing REAL, USEFUL feedback:\n"
        "1. CRITICAL RISK ASSESSMENT: Identify projects at risk based on status, priority, and lack of recent updates.\n"
        "2. STRATEGIC ALIGNMENT: Suggest how these projects align with modern SAP/Cloud best practices.\n"
        "3. ACCELERATION OPPORTUNITIES: Where can we move faster? What blockers can be removed?\n"
        "4. TECHNICAL ADVICE: Use the Document Summaries to give specific technical recommendations.\n\n"
        "IMPORTANT: Output the report as valid HTML code with inline CSS for styling. "
        "Use a 'Cyberpunk' aesthetic: dark background is already provided by the container, so use transparent backgrounds. "
        "Use colors like #00f3ff (Cyan), #ff00ff (Magenta), #fcee0a (Yellow), and #ff2a2a (Red). "
        "Use font-family: 'Rajdhani', sans-serif for headers and 'Roboto Mono', monospace for text. "
        "Format with <h2>, <h3>, <ul>, <li>, <p> tags. "
        "Do NOT include <html>, <head>, or <body> tags, just the content div."
    )

    try:
        completion = client.chat.completions.create(
            model=DEFAULT_MODEL,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"Here is the project portfolio:\n\n{project_context}"}
            ],
            max_tokens=2000,
        )
        content = completion.choices[0].message.content
        # Strip markdown code blocks if present
        if content.startswith("```html"):
            content = content[7:]
        if content.startswith("```"):
            content = content[3:]
        if content.endswith("```"):
            content = content[:-3]
        return content
    except Exception as e:
        _log(f"Error generating AMD AI report: {e}")
        return f"Error generating report: {str(e)}"


@app.route("/api/dems/report/ai", methods=["POST"])
def amd_ai_report():
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "Unauthorized"}), 401

    data = request.get_json() or {}
    projects = data.get("projects", [])

    if not projects:
        return jsonify({"error": "No projects provided for analysis."}), 400

    # Filter out SAP related projects
    filtered_projects = []
    for p in projects:
        # Check name, title, and notes for "SAP"
        text_content = (
            (p.get("name") or "") + 
            (p.get("title") or "") + 
            " ".join([n.get("text", "") if isinstance(n, dict) else str(n) for n in p.get("notes", [])])
        ).lower()
        
        if "sap" not in text_content:
            filtered_projects.append(p)

    if not filtered_projects:
        return jsonify({"error": "No eligible projects found (SAP projects are excluded)."}), 400

    report_text = generate_amd_ai_report_logic(filtered_projects)
    return jsonify({"report": report_text})



def generate_ai_solution_analysis_logic(project):
    """
    Generates a comprehensive solution analysis for a single project.
    """
    name = project.get("name", "Unknown")
    title = project.get("title", "No Title")
    doc_summary = project.get("doc_summary", "No document summary available.")
    notes = project.get("notes", [])
    
    notes_text = "\n".join([n.get("text", "") if isinstance(n, dict) else str(n) for n in notes][-5:])

    system_prompt = (
        "You are an Expert Solution Architect. "
        "Analyze the following project request and provide a comprehensive solution analysis.\n"
        "Include:\n"
        "1. Problem Statement Analysis\n"
        "2. Proposed Solution Architecture (High Level)\n"
        "3. Key Technical Components (SAP modules, Cloud services, etc.)\n"
        "4. Implementation Steps & Risks\n\n"
        "Use the Document Summary as the primary source of requirements.\n\n"
        "IMPORTANT: Output the report as valid HTML code with inline CSS for styling. "
        "Use a 'Cyberpunk' aesthetic: dark background is already provided by the container, so use transparent backgrounds. "
        "Use colors like #00f3ff (Cyan), #ff00ff (Magenta), #fcee0a (Yellow), and #ff2a2a (Red). "
        "Use font-family: 'Rajdhani', sans-serif for headers and 'Roboto Mono', monospace for text. "
        "Format with <h2>, <h3>, <ul>, <li>, <p> tags. "
        "Do NOT include <html>, <head>, or <body> tags, just the content div."
    )
    
    user_content = (
        f"Project: {name} - {title}\n"
        f"Document Summary: {doc_summary}\n"
        f"Recent Notes: {notes_text}\n"
    )

    try:
        completion = client.chat.completions.create(
            model=DEFAULT_MODEL,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_content}
            ],
            max_tokens=1500,
        )
        return completion.choices[0].message.content
    except Exception as e:
        return f"Error generating analysis: {str(e)}"

@app.route("/api/dems/project/<project_id>/analysis", methods=["POST"])
def generate_ai_solution_analysis(project_id):
    maybe = require_auth()
    if maybe: return maybe
    
    projects = load_dems()
    project = next((p for p in projects if p.get("id") == project_id), None)
    
    if not project:
        return jsonify({"error": "Project not found"}), 404
        
    analysis = generate_ai_solution_analysis_logic(project)
    return jsonify({"analysis": analysis})

@app.route("/api/files", methods=["GET"])
def list_files():
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "Unauthorized"}), 401
    
    try:
        files = []
        if os.path.exists(UPLOAD_FOLDER):
            for f in os.listdir(UPLOAD_FOLDER):
                path = os.path.join(UPLOAD_FOLDER, f)
                if os.path.isfile(path):
                    stat = os.stat(path)
                    size_mb = round(stat.st_size / (1024 * 1024), 2)
                    date_str = datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M')
                    files.append({
                        "name": f,
                        "size": f"{size_mb} MB",
                        "date": date_str
                    })
        return jsonify({"files": files})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/files/<filename>", methods=["GET"])
def download_file(filename):
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "Unauthorized"}), 401
    
    try:
        return send_from_directory(UPLOAD_FOLDER, filename, as_attachment=True)
    except Exception as e:
        return jsonify({"error": str(e)}), 404

@app.route("/api/files/<filename>/delete", methods=["POST"])
def delete_file(filename):
    maybe = require_auth()
    if maybe is not None:
        return jsonify({"error": "Unauthorized"}), 401
    
    try:
        path = os.path.join(UPLOAD_FOLDER, filename)
        if os.path.exists(path):
            os.remove(path)
            return jsonify({"success": True})
        else:
            return jsonify({"error": "File not found"}), 404
    except Exception as e:
        return jsonify({"error": str(e)}), 500


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=8080)
